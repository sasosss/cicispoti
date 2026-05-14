#!/usr/bin/env python3
"""
Genera una presentazione PowerPoint sulla provincia di Benevento.
Richiede: python-pptx, Pillow, requests
"""

from __future__ import annotations

import sys
import time
from io import BytesIO
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

# Palette: tramonto, vino, terra, paglia
WINE = RGBColor(0x6B, 0x1F, 0x1F)
BURGUNDY = RGBColor(0x8B, 0x2E, 0x2E)
TERRA = RGBColor(0x5C, 0x3A, 0x21)
SUNSET = RGBColor(0xD4, 0x8A, 0x4A)
GOLD = RGBColor(0xE8, 0xC4, 0x7A)
CREAM = RGBColor(0xFA, 0xF0, 0xE4)
DEEP = RGBColor(0x2A, 0x18, 0x12)
ASSETS_DIR = Path(__file__).resolve().parent.parent / "presentations" / "benevento_assets"
OUTPUT_PATH = Path(__file__).resolve().parent.parent / "presentations" / "Benevento_terra_storia_vino_sapori.pptx"

# Immagini libere (Wikimedia Commons) — uso thumb per dimensioni gestibili
IMAGE_URLS = {
    "cover_arch": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/6/6e/"
        "Arco_di_Traiano_%28Benevento%2C_Campania%29.jpg/"
        "1280px-Arco_di_Traiano_%28Benevento%2C_Campania%29.jpg"
    ),
    "hills": "https://upload.wikimedia.org/wikipedia/commons/5/54/Vineyard_in_Tuscany.jpg",
    "campania_map": "https://upload.wikimedia.org/wikipedia/commons/5/53/Campania_Provinces_ar.png",
    "wine_glass": "https://upload.wikimedia.org/wikipedia/commons/e/e3/Wine_being_poured_into_a_glass.jpg",
    "vineyard_rows": (
        "https://upload.wikimedia.org/wikipedia/commons/e/ef/"
        "Grapevines%2C_Llanerch_Vineyard%2C_Vale_of_Glamorgan_-_geograph.org.uk_-_1079208.jpg"
    ),
    "cavatelli": "https://upload.wikimedia.org/wikipedia/commons/d/d6/Cavatelli.jpg",
    "ragu": (
        "https://upload.wikimedia.org/wikipedia/commons/7/75/"
        "Liat_Portal_for_Foodie_Disorder_-_Homemade_Bolognese_Ragu_Sauce.jpg"
    ),
}


def download_image(url: str, stem: Path, timeout: int = 45) -> bool:
    """Scarica e normalizza in JPEG (PowerPoint non accetta alcuni formati es. MPO)."""
    stem.parent.mkdir(parents=True, exist_ok=True)
    out = stem.with_suffix(".jpg")
    for attempt in range(4):
        try:
            r = requests.get(
                url,
                timeout=timeout,
                headers={"User-Agent": "Mozilla/5.0 (compatible; BeneventoDeckBuilder/1.0)"},
            )
            if r.status_code == 429:
                time.sleep(2.0 * (attempt + 1))
                continue
            r.raise_for_status()
            raw = r.content
            im = Image.open(BytesIO(raw))
            if im.mode in ("RGBA", "P", "LA"):
                im = im.convert("RGBA")
                bg = Image.new("RGB", im.size, (255, 255, 255))
                bg.paste(im, mask=im.split()[-1])
                im = bg
            else:
                im = im.convert("RGB")
            im.save(out, "JPEG", quality=90, optimize=True)
            return True
        except Exception as e:
            if attempt == 3:
                print(f"WARN: impossibile scaricare {url}: {e}", file=sys.stderr)
                return False
            time.sleep(1.5 * (attempt + 1))
    return False


def fetch_all_assets() -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}
    for key, url in IMAGE_URLS.items():
        stem = ASSETS_DIR / key
        if download_image(url, stem):
            paths[key] = stem.with_suffix(".jpg")
        time.sleep(0.75)
    return paths


def _set_slide_transition(slide, fade: bool = True) -> None:
    """Aggiunge transizione fade alla slide (OOXML)."""
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    if fade:
        xml = (
            f'<p:transition xmlns:p="{ns}" spd="med" advTm="0">'
            f'<p:fade/></p:transition>'
        )
    else:
        xml = (
            f'<p:transition xmlns:p="{ns}" spd="med" advTm="0">'
            f'<p:push dir="u"/></p:transition>'
        )
    el = parse_xml(xml)
    slide._element.append(el)


def add_gradient_band(slide, top: bool, color1: RGBColor, color2: RGBColor) -> None:
    h = Inches(0.35)
    y = Inches(0) if top else Inches(7.15)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, Inches(13.33), h
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    shape.line.fill.background()


def style_title(tf, size_pt: float = 44, color: RGBColor = CREAM, font: str = "Georgia") -> None:
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.bold = True
            run.font.color.rgb = color
            run.font.small_caps = True


def slide_blank(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_full_bleed_image(slide, path: Path, brightness: float = 0.55) -> None:
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(0), width=Inches(13.33))
    target_h = int(Inches(7.5))
    old_w, old_h = pic.width, pic.height
    pic.height = target_h
    pic.width = int(old_w * target_h / old_h)
    pic.left = int((Inches(13.33) - pic.width) / 2)
    pic.top = 0
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(13.33)), int(Inches(7.5)))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x1A, 0x0E, 0x0C)
    overlay.fill.transparency = 1.0 - brightness
    overlay.line.fill.background()


def build_presentation(asset_paths: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = int(Inches(13.33))
    prs.slide_height = int(Inches(7.5))

    # --- Copertina ---
    s = slide_blank(prs)
    if "cover_arch" in asset_paths:
        add_full_bleed_image(s, asset_paths["cover_arch"], 0.5)
    else:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(13.33)), int(Inches(7.5)))
        bg.fill.gradient()
        bg.fill.gradient_angle = 45.0
        bg.fill.gradient_stops[0].color.rgb = BURGUNDY
        bg.fill.gradient_stops[1].color.rgb = TERRA
        bg.line.fill.background()

    add_gradient_band(s, True, SUNSET, BURGUNDY)
    add_gradient_band(s, False, TERRA, DEEP)

    title_box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.4))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Benevento: terra di storia, vino e sapori"
    style_title(tf, 40, CREAM, "Georgia")
    sub = s.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.6))
    stf = sub.text_frame
    stf.paragraphs[0].text = "◆  Un viaggio tra colline, cantine e tavole campane  ◆"
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for r in stf.paragraphs[0].runs:
        r.font.name = "Calibri Light"
        r.font.size = Pt(20)
        r.font.color.rgb = GOLD
    _set_slide_transition(s)

    # --- Intro: storia e territorio (layout 50/50, testo leggibile) ---
    s = slide_blank(prs)
    if "hills" in asset_paths:
        s.shapes.add_picture(str(asset_paths["hills"]), Inches(0), Inches(0), height=Inches(7.5))
        ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(6.45)), int(Inches(7.5)))
        ov.fill.solid()
        ov.fill.fore_color.rgb = RGBColor(0x2A, 0x12, 0x10)
        ov.fill.transparency = 0.35
        ov.line.fill.background()
    else:
        leftg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(6.45)), int(Inches(7.5)))
        leftg.fill.gradient()
        leftg.fill.gradient_angle = 0.0
        leftg.fill.gradient_stops[0].color.rgb = BURGUNDY
        leftg.fill.gradient_stops[1].color.rgb = SUNSET
        leftg.line.fill.background()

    right = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(Inches(6.45)), 0, int(Inches(6.88)), int(Inches(7.5))
    )
    right.fill.solid()
    right.fill.fore_color.rgb = CREAM
    right.line.fill.background()

    tb = s.shapes.add_textbox(Inches(6.75), Inches(0.45), Inches(6.35), Inches(6.65))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "La provincia di Benevento"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.color.rgb = BURGUNDY

    intro_body = (
        "Nel cuore della Campania interna, la provincia di Benevento custodisce millenni di storia "
        "e colline morbide punteggiate da uliveti e vigneti: un paesaggio che al tramonto si accende "
        "di toni ambra e porpora.\n\n"
        "Città sannita e poi crocevia romano, Benevento conserva capolavori come l’Arco di Traiano "
        "e il complesso longobardo di Santa Sofia (Patrimonio UNESCO). Nelle contrade, il tempo "
        "si misura ancora a passi di processione, a calici alzati in cantina e a tavolate che profumano "
        "di ragù lungo cottura."
    )
    p2 = tf.add_paragraph()
    p2.text = intro_body
    p2.font.size = Pt(13.5)
    p2.font.name = "Calibri"
    p2.font.color.rgb = DEEP
    p2.line_spacing = 1.22

    p3 = tf.add_paragraph()
    p3.text = (
        "◆  Storia e cultura: dalle origini sannitiche alla romanità, dal Medioevo longobardo "
        "alle pieghe del Sannio contemporaneo — la memoria qui è paesaggio, e il paesaggio è accoglienza."
    )
    p3.font.size = Pt(13.5)
    p3.font.name = "Calibri"
    p3.font.color.rgb = BURGUNDY
    p3.font.bold = True
    p3.line_spacing = 1.2

    p4 = tf.add_paragraph()
    p4.text = (
        "◆  Paesaggi collinari: successioni morbide di ulivi e vigneti, borghi arroccati, "
        "strade bianche profumate di erba secca dopo la pioggia, e silenzi punteggiati solo dal vento."
    )
    p4.font.size = Pt(13.5)
    p4.font.name = "Calibri"
    p4.font.color.rgb = DEEP
    p4.line_spacing = 1.2

    p5 = tf.add_paragraph()
    p5.text = (
        "◆  Tradizioni popolari: feste patronali, tarantelle e canti di lavoro, mercati contadini "
        "e ricami di pizzi — tutto ciò che ancora oggi «fa comunità» attorno al tavolo e alla piazza."
    )
    p5.font.size = Pt(13.5)
    p5.font.name = "Calibri"
    p5.font.color.rgb = DEEP
    p5.line_spacing = 1.2
    _set_slide_transition(s)

    # --- Leggende e streghe ---
    s = slide_blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(13.33)), int(Inches(0.55)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP
    bar.line.fill.background()

    if "campania_map" in asset_paths:
        s.shapes.add_picture(str(asset_paths["campania_map"]), Inches(0.35), Inches(0.85), width=Inches(4.6))
    map_cap = s.shapes.add_textbox(Inches(0.35), Inches(5.35), Inches(4.6), Inches(0.35))
    map_cap.text_frame.paragraphs[0].text = "Mappa stilizzata — Campania e Benevento"
    map_cap.text_frame.paragraphs[0].font.size = Pt(10)
    map_cap.text_frame.paragraphs[0].font.italic = True
    map_cap.text_frame.paragraphs[0].font.color.rgb = TERRA

    tb = s.shapes.add_textbox(Inches(5.2), Inches(0.75), Inches(7.7), Inches(6.5))
    tf = tb.text_frame
    tf.paragraphs[0].text = "Curiosità: le streghe di Benevento"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BURGUNDY

    legend_text = (
        "Fin dal Medioevo, cronache e leggende hanno associato Benevento a raduni notturni e "
        "riti «oltre» la norma cristiana: si parlava di «noce di Benevento» come luogo simbolico "
        "d’incontro, e di presenze femminili legate a saperi erboristici, medicina popolare e "
        "figure emarginate dalla società del tempo.\n\n"
        "Oggi quelle narrazioni — filtrate dalla storia e dalla letteratura — diventano patrimonio "
        "evocativo: non un manuale di magia, ma un invito a esplorare il mistero delle colline, "
        "dei boschi e delle voci tramandate, tra folklore e ironia sannita.\n\n"
        "◆  Il filo conduttore è sempre lo stesso: una terra che sa raccontare storie al fuoco del camino "
        "e al calice di un vino corposo."
    )
    p = tf.add_paragraph()
    p.text = legend_text
    p.font.size = Pt(15)
    p.font.name = "Calibri"
    p.font.color.rgb = DEEP
    p.line_spacing = 1.2
    _set_slide_transition(s, fade=False)

    # --- Vino: Aglianico del Taburno (hero) ---
    s = slide_blank(prs)
    if "wine_glass" in asset_paths:
        s.shapes.add_picture(str(asset_paths["wine_glass"]), Inches(7.0), Inches(0.35), width=Inches(5.9))
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(6.55)), int(Inches(7.5)))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x3D, 0x1A, 0x1A)
    panel.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.45), Inches(5.75), Inches(6.6))
    tf = tb.text_frame
    tf.paragraphs[0].text = "Aglianico del Taburno"
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD

    wine_txt = (
        "Origine e territorio\n"
        "L’Aglianico del Taburno nasce sulle pendici del massiccio vulcanico del Taburno–Camposauro, "
        "tra suoli ricchi di minerali e brezze che favoriscono lente maturazioni. Il disciplinare "
        "valorizza un legame stretto tra varietà autoctona, altitudine e microclima collinare.\n\n"
        "Carattere in cantina\n"
        "Struttura importante, tannini presenti ma nobili nel tempo, grande acidità che sostiene "
        "l’invecchiamento: è un vino da «aspettare», capace di raccontare annate e cantine con "
        "eleganza quasi cinematografica.\n\n"
        "Sensorialità e tavola\n"
        "Colore rubino intenso con riflessi granati dopo gli anni. Profumi di frutta nera matura, "
        "spezie dolci, tabacco, sottobosco e, spesso, note vulcaniche. In abbinamento: carni rosse "
        "alla brace, selvaggina, funghi porcini, formaggi stagionati e — naturalmente — i ragù "
        "lenti della tradizione beneventana."
    )
    p = tf.add_paragraph()
    p.text = wine_txt
    p.font.size = Pt(13.5)
    p.font.name = "Calibri"
    p.font.color.rgb = CREAM
    p.line_spacing = 1.18
    _set_slide_transition(s)

    # --- Vigneti panoramici ---
    s = slide_blank(prs)
    if "vineyard_rows" in asset_paths:
        pic = s.shapes.add_picture(str(asset_paths["vineyard_rows"]), 0, 0, width=int(Inches(13.33)))
        target_h = int(Inches(7.5))
        old_w, old_h = pic.width, pic.height
        pic.height = target_h
        pic.width = int(old_w * target_h / old_h)
        pic.top = 0
        pic.left = int((Inches(13.33) - pic.width) / 2)
    scrim = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(13.33)), int(Inches(7.5)))
    scrim.fill.solid()
    scrim.fill.fore_color.rgb = DEEP
    scrim.fill.transparency = 0.45
    scrim.line.fill.background()

    cap = s.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.0))
    ctf = cap.text_frame
    ctf.paragraphs[0].text = "◆  Filari al tramonto  —  l’orchestra silenziosa del Taburno"
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.paragraphs[0].font.size = Pt(26)
    ctf.paragraphs[0].font.name = "Georgia"
    ctf.paragraphs[0].font.color.rgb = GOLD
    ctf.paragraphs[0].font.bold = True

    foot = s.shapes.add_textbox(Inches(1.2), Inches(5.85), Inches(10.9), Inches(1.2))
    ftf = foot.text_frame
    ftf.paragraphs[0].text = (
        "Camminare tra i vigneti significa respirare terra bollente, sentire il fruscio delle foglie "
        "e intuire, nell’aria, il profumo che diventerà — mesi dopo — bouquet in bicchiere. "
        "È il capitolo visivo di un mini-documentario che profuma di mosto e di sole al tramonto."
    )
    ftf.paragraphs[0].font.size = Pt(15)
    ftf.paragraphs[0].font.name = "Calibri"
    ftf.paragraphs[0].font.color.rgb = CREAM
    ftf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_slide_transition(s)

    # --- Cavatelli: ingredienti e tradizione ---
    s = slide_blank(prs)
    if "cavatelli" in asset_paths:
        s.shapes.add_picture(str(asset_paths["cavatelli"]), Inches(0.35), Inches(1.05), width=Inches(5.85))
    if "ragu" in asset_paths:
        s.shapes.add_picture(str(asset_paths["ragu"]), Inches(6.75), Inches(1.05), width=Inches(6.0))

    title = s.shapes.add_textbox(Inches(0.35), Inches(0.35), Inches(12.5), Inches(0.55))
    ttf = title.text_frame
    ttf.paragraphs[0].text = "Cavatelli al ragù beneventano — identità di collina"
    ttf.paragraphs[0].font.size = Pt(26)
    ttf.paragraphs[0].font.name = "Georgia"
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = BURGUNDY

    tb = s.shapes.add_textbox(Inches(0.35), Inches(4.55), Inches(12.6), Inches(2.75))
    tf = tb.text_frame
    tf.paragraphs[0].text = (
        "Ingredienti tradizionali\n"
        "◆  Cavatelli (pasta fresca di semola e acqua, «strascinati» sul tagliere)\n"
        "◆  Ragù di carne mista (spesso manzo e maiale) con odori classici: sedano, carota, cipolla\n"
        "◆  Passata o polpa di pomodoro maturo, vino rosso, olio extravergine\n"
        "◆  Pecorino stagionato o ricotta salata, a seconda della famiglia e della festa\n\n"
        "Cucina contadina\n"
        "Il piatto nasce dalla lentezza: fuoco basso, pentola che «canta» per ore, profumi che "
        "attraversano la casa. È economia intelligente — trasformare tagli umili in sapore profondo — "
        "e convivialità: la tavola allungata, il pane che fa la scarpetta, il vino che accompagna "
        "senza coprire.\n\n"
        "◆  Accanto: cavatelli fatti in casa (foto) e ragù lento in cottura (immagine di riferimento "
        "da cucina italiana) — invito visivo ad assaporare Benevento anche con gli occhi."
    )
    for pr in tf.paragraphs:
        pr.font.size = Pt(13)
        pr.font.name = "Calibri"
        pr.font.color.rgb = DEEP
        pr.line_spacing = 1.15
    _set_slide_transition(s)

    # --- Chiusura emozionale ---
    s = slide_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(Inches(13.33)), int(Inches(7.5)))
    bg.fill.gradient()
    bg.fill.gradient_angle = 135.0
    bg.fill.gradient_stops[0].color.rgb = BURGUNDY
    bg.fill.gradient_stops[1].color.rgb = SUNSET
    bg.line.fill.background()

    deco = s.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5))
    dtf = deco.text_frame
    dtf.paragraphs[0].text = "◆     ◆     ◆"
    dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    dtf.paragraphs[0].font.size = Pt(28)
    dtf.paragraphs[0].font.color.rgb = GOLD

    closing = s.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.4))
    ctf = closing.text_frame
    ctf.paragraphs[0].text = "Benevento non si visita soltanto… si assapora."
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.paragraphs[0].font.size = Pt(36)
    ctf.paragraphs[0].font.name = "Georgia"
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = CREAM
    ctf.paragraphs[0].line_spacing = 1.1

    tag = s.shapes.add_textbox(Inches(1.0), Inches(4.35), Inches(11.3), Inches(1.8))
    tgf = tag.text_frame
    tgf.paragraphs[0].text = (
        "Tra pietre antiche e filari moderni, tra cantine silenziose e cucine rumorose, "
        "la provincia di Benevento offre un racconto continuo: storico, umano, gastronomico.\n\n"
        "Chi arriva qui porta via non solo fotografie panoramiche, ma anche il sapore di un luogo "
        "che sa fermare l’orologio — almeno per un bicchiere, almeno per un piatto."
    )
    tgf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tgf.paragraphs[0].font.size = Pt(16)
    tgf.paragraphs[0].font.name = "Calibri"
    tgf.paragraphs[0].font.color.rgb = CREAM
    _set_slide_transition(s)

    credits = (
        "Immagini da Wikimedia Commons (licenze libere): "
        "Arco di Traiano (Benevento) — Anna Eden; "
        "Vigneto — Markus Bärlocher; "
        "Mappa province Campania — Ali1; "
        "Vino versato in calice — Flickr/Vinitaly; "
        "Filari — John Lord / geograph.org.uk; "
        "Cavatelli — Luigi Scorcia; "
        "Ragù in cottura — Liat Portal. "
        "Rigenera con: python3 scripts/build_benevento_presentation.py\n\n"
        "Animazioni (PowerPoint): Transizioni → Dissolvenza o Spostamento, 1,0–1,25 s; "
        "Animazioni → titoli con Dissolvenza entrata / Scorrimento da basso; "
        "tempistiche per effetto documentario."
    )
    try:
        prs.slides[0].notes_slide.notes_text_frame.text = credits
    except Exception:
        pass

    return prs


def main() -> int:
    asset_paths = fetch_all_assets()
    if len(asset_paths) < 3:
        print("Alcune immagini non sono state scaricate; la presentazione userà sfondi gradiente.", file=sys.stderr)
    prs = build_presentation(asset_paths)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))

    print(f"Salvato: {OUTPUT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
